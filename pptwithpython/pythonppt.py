import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

def add_text_to_slide(slide, text, font_name="Garamond", font_size=28, line_spacing=1.5):
    """Adds text to a slide with proper formatting and alignment."""
    left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = text
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    p.line_spacing = line_spacing

def split_text_into_chunks(text, max_chars=500):
    """Splits long text into smaller chunks for separate slides."""
    words = text.split()
    slides = []
    chunk = ""
    
    for word in words:
        if len(chunk) + len(word) + 1 <= max_chars:
            chunk += " " + word
        else:
            slides.append(chunk.strip())
            chunk = word
    
    if chunk:
        slides.append(chunk.strip())
    
    return slides

# Create PowerPoint presentation
presentation = Presentation()

# Read content from file
txt_file = "sample_slide1_input.txt"

with open(txt_file, "r") as file:
    content = file.read()

# Split content into multiple slides
text_chunks = split_text_into_chunks(content, max_chars=500)

# Add slides dynamically
for chunk in text_chunks:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_text_to_slide(slide, chunk)

# Save the presentation
ppt_file = "auto_slides_presentation.pptx"
presentation.save(ppt_file)

print(f"Presentation saved as {ppt_file}")
